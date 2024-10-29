# app/utils.py
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import sys
import requests
import json
import uuid
import logging
from datetime import datetime, timedelta
from .models import Message, LearningCompanion
import random
import string

ALLOWED_EXTENSIONS = {'pdf', 'docx', 'doc', 'ppt', 'pptx'}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def generate_conversation_id():
    return str(uuid.uuid4())

def convert_file_to_text(filepath, filename):
    extension = filename.rsplit('.', 1)[1].lower()
    if extension == 'pdf':
        return convert_pdf_to_text(filepath)
    elif extension in ['docx', 'doc']:
        return convert_docx_to_text(filepath)
    elif extension in ['ppt', 'pptx']:
        return convert_ppt_to_text(filepath)
    else:
        return "Unsupported file format."

# Convert PDF to text
def convert_pdf_to_text(filepath):
    text = ''
    try:
        with open(filepath, 'rb') as f:
            reader = PdfReader(f)
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text += page.extract_text()
    except Exception as e:
        return f"Error extracting text from PDF: {str(e)}"
    return text

# Convert DOCX to text
def convert_docx_to_text(filepath):
    text = ''
    try:
        doc = Document(filepath)
        for para in doc.paragraphs:
            text += para.text + '\n'
    except Exception as e:
        return f"Error extracting text from DOCX: {str(e)}"
    return text

# Convert PPT/PPTX to text
def convert_ppt_to_text(filepath):
    text = ''
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
    except Exception as e:
        return f"Error extracting text from PPT/PPTX: {str(e)}"
    return text

def check_message_limit(user_id, companion_id, limit=45):
    """Checks if a student has exceeded their message limit."""
    one_hour_ago = datetime.utcnow() - timedelta(hours=1)
    message_count = Message.query.filter_by(user_id=user_id, companion_id=companion_id).filter(Message.timestamp >= one_hour_ago).count()
    return message_count < limit

# app/utils.py

import sys
import os
import requests
import json
import uuid
import logging

class Coze:
    def __init__(self,
                 bot_id=None,
                 api_token=None,
                 user_id="default_user",
                 conversation_id=None,
                 stream=False):
        self.bot_id = bot_id
        self.api_token = api_token
        self.user_id = user_id
        self.conversation_id = conversation_id or self.generate_conversation_id()
        self.stream = stream
        self.url = 'https://api.coze.com/open_api/v2/chat'
        self.headers = {
            'Authorization': f'Bearer {self.api_token}',
            'Content-Type': 'application/json',
            'Accept': '*/*',
            'Connection': 'keep-alive'
        }
        logging.basicConfig(stream=sys.stderr, level=logging.INFO)
        self.logger = logging.getLogger(__name__)

    @staticmethod
    def generate_conversation_id():
        return str(uuid.uuid4())

    @staticmethod
    def build_messages(history=None):
        """
        Builds the message list to send to the Coze API.
        :param history: A list of (content, is_user) tuples.
        :return: A list of dictionaries formatted for the Coze API.
        """
        messages = []
        if history:
            for content, is_user in history:
                role = 'user' if is_user else 'assistant'
                messages.append({
                    "role": role,
                    "content": content
                })
        return messages

    def chat(self, query, history=None):
        """
        Sends a message to the Coze API and streams the assistant's response incrementally.
        :param query: The user's message.
        :param history: A list of (content, is_user) tuples representing the chat history.
        :return: A generator that yields the assistant's response in chunks.
        """
        payload = {
            "bot_id": self.bot_id,
            "user": self.user_id,
            "query": query,
            "stream": self.stream
        }

        # Include conversation_id if provided
        if self.conversation_id:
            payload["conversation_id"] = self.conversation_id

        # Include chat_history if provided
        if history:
            payload["chat_history"] = self.build_messages(history)

        try:
            response = requests.post(self.url, headers=self.headers, json=payload, stream=self.stream)
            response.raise_for_status()
        except requests.RequestException as e:
            self.logger.error(f"Request error: {e}")
            yield "An error occurred while processing your request."
            return

        if not self.stream:
            try:
                data = response.json()
                yield self.get_response(data)
            except json.JSONDecodeError as e:
                self.logger.error(f"JSON decode error: {e}")
                yield "Invalid response format received."
        else:
            for line in response.iter_lines():
                if line:
                    try:
                        decoded_line = line.decode('utf-8').strip()
                        if not decoded_line.startswith('data:'):
                            continue
                        json_str = decoded_line.split("data:")[-1].strip()
                        if not json_str:
                            continue
                        data = json.loads(json_str)
                        event = data.get('event')
                        if event == 'done':
                            break
                        elif event == 'message':
                            message = data.get('message', {})
                            msg_type = message.get('type')
                            if msg_type == 'answer':
                                content = message.get('content', '')
                                yield content
                    except json.JSONDecodeError as e:
                        self.logger.error(f"JSON decode error: {e}")
                        continue
            # Optionally, yield a completion signal
            yield "[DONE]"

        response.close()

    def get_response(self, data):
        """
        Processes the API response to extract the assistant's response.
        :param data: The API response data.
        :return: The response content from the assistant.
        """
        if 'messages' in data:
            messages = data['messages']
            response_content = ''
            for msg in messages:
                if msg['role'] == 'assistant' and msg['type'] == 'answer':
                    response_content += msg['content']
            return response_content
        else:
            self.logger.error(f"No messages in response: {data}")
            return "No messages received from the server."

    def __call__(self, query, history=None):
        """
        Enables the Coze instance to be called like a function.
        :param query: The user's message.
        :param history: The conversation history as a list of tuples (content, is_user).
        :return: A generator that yields the assistant's response in chunks.
        """
        return self.chat(query, history=history)

    def reset_conversation(self):
        """
        Resets the conversation by generating a new conversation_id.
        """
        self.conversation_id = self.generate_conversation_id()
        self.logger.info("Conversation has been reset with new conversation_id.")

